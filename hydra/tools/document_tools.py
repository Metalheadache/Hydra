"""
Document generation and reading tools: DOCX, XLSX, PPTX, PDF.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import TYPE_CHECKING, Any

import structlog

from hydra.models import ToolResult
from hydra.tools.base import BaseTool
from hydra.tools.file_tools import _ensure_output_dir, _safe_filepath, _DEFAULT_OUTPUT_DIR

if TYPE_CHECKING:
    from hydra.state_manager import StateManager

logger = structlog.get_logger(__name__)


# ── WriteDocxTool ─────────────────────────────────────────────────────────────

class WriteDocxTool(BaseTool):
    """Write a professional Word document (.docx) from markdown-like content."""

    name = "write_docx"
    description = (
        "Create a Word document (.docx) from markdown-like text. "
        "Supports headings (# / ##), bullet lists (- item), bold (**text**), italic (*text*), "
        "and plain paragraphs. Returns the filepath of the created file."
    )
    parameters = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Output filename (e.g. 'report.docx'). Extension added if missing.",
            },
            "content": {
                "type": "string",
                "description": "Markdown-like content to write into the document.",
            },
            "title": {
                "type": "string",
                "description": "Optional document title added as a Title-style paragraph at the top.",
            },
            "output_dir": {
                "type": "string",
                "description": f"Directory to write the file. Defaults to '{_DEFAULT_OUTPUT_DIR}'.",
            },
        },
        "required": ["filename", "content"],
    }

    def __init__(
        self,
        output_dir: str = _DEFAULT_OUTPUT_DIR,
        state_manager: "StateManager | None" = None,
    ) -> None:
        self._output_dir = output_dir
        self._state_manager = state_manager

    async def execute(
        self,
        filename: str,
        content: str,
        title: str | None = None,
        output_dir: str | None = None,
    ) -> ToolResult:
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            return ToolResult(success=False, error="python-docx is not installed. Run: pip install python-docx")

        try:
            effective_dir = output_dir if output_dir is not None else self._output_dir
            output_path = _ensure_output_dir(effective_dir)
            if not filename.endswith(".docx"):
                filename += ".docx"
            filepath = _safe_filepath(output_path, filename)
            if filepath is None:
                return ToolResult(success=False, error="Path traversal blocked")

            doc = Document()

            # ── Styling tweaks ──────────────────────────────────────────────
            style = doc.styles["Normal"]
            style.font.name = "Calibri"
            style.font.size = Pt(11)

            if title:
                t = doc.add_paragraph(title, style="Title")
                t.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # ── Parse content ───────────────────────────────────────────────
            for line in content.splitlines():
                stripped = line.strip()

                if stripped.startswith("### "):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith("# "):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith("- ") or stripped.startswith("* "):
                    # Bullet list item
                    para = doc.add_paragraph(style="List Bullet")
                    _add_inline_formatting(para, stripped[2:])
                elif stripped == "":
                    # Blank line → just continue (paragraph breaks handled naturally)
                    pass
                else:
                    para = doc.add_paragraph()
                    _add_inline_formatting(para, stripped)

            doc.save(str(filepath))

            if self._state_manager is not None:
                await self._state_manager.register_file(filename, str(filepath))

            logger.info("docx_written", filepath=str(filepath))
            return ToolResult(success=True, data={"filepath": str(filepath)})

        except Exception as exc:
            logger.error("write_docx_failed", error=str(exc))
            return ToolResult(success=False, error=f"Failed to write DOCX: {exc}")


def _add_inline_formatting(para, text: str) -> None:
    """Add a run to a paragraph with inline bold/italic parsing."""
    # Pattern: **bold**, *italic*
    pattern = re.compile(r"(\*\*(.+?)\*\*|\*(.+?)\*|([^*]+))")
    for match in pattern.finditer(text):
        if match.group(2):  # **bold**
            run = para.add_run(match.group(2))
            run.bold = True
        elif match.group(3):  # *italic*
            run = para.add_run(match.group(3))
            run.italic = True
        elif match.group(4):  # plain text
            para.add_run(match.group(4))


# ── WriteXlsxTool ─────────────────────────────────────────────────────────────

class WriteXlsxTool(BaseTool):
    """Write data to an Excel workbook (.xlsx)."""

    name = "write_xlsx"
    description = (
        "Create an Excel workbook (.xlsx). Accepts either: "
        "(a) `sheets` — a list of {name, headers, rows} for multi-sheet workbooks, or "
        "(b) `data` — a simple list of dicts for a single sheet. "
        "Features: auto-column width, bold headers, frozen top row, auto-filter. "
        "Returns the filepath."
    )
    parameters = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Output filename (e.g. 'data.xlsx').",
            },
            "sheets": {
                "type": "array",
                "description": "List of sheet definitions: [{name, headers, rows}].",
                "items": {
                    "type": "object",
                    "properties": {
                        "name": {"type": "string"},
                        "headers": {"type": "array", "items": {"type": "string"}},
                        "rows": {"type": "array", "items": {"type": "array"}},
                    },
                    "required": ["headers", "rows"],
                },
            },
            "data": {
                "type": "array",
                "description": "Simple list of dicts for a single sheet (alternative to `sheets`).",
                "items": {"type": "object"},
            },
            "output_dir": {
                "type": "string",
                "description": f"Directory to write the file. Defaults to '{_DEFAULT_OUTPUT_DIR}'.",
            },
        },
        "required": ["filename"],
    }

    def __init__(
        self,
        output_dir: str = _DEFAULT_OUTPUT_DIR,
        state_manager: "StateManager | None" = None,
    ) -> None:
        self._output_dir = output_dir
        self._state_manager = state_manager

    async def execute(
        self,
        filename: str,
        sheets: list[dict] | None = None,
        data: list[dict] | None = None,
        output_dir: str | None = None,
    ) -> ToolResult:
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.utils import get_column_letter
        except ImportError:
            return ToolResult(success=False, error="openpyxl is not installed. Run: pip install openpyxl")

        try:
            effective_dir = output_dir if output_dir is not None else self._output_dir
            output_path = _ensure_output_dir(effective_dir)
            if not filename.endswith(".xlsx"):
                filename += ".xlsx"
            filepath = _safe_filepath(output_path, filename)
            if filepath is None:
                return ToolResult(success=False, error="Path traversal blocked")

            # Build sheet definitions
            if sheets is None and data is not None:
                # Convert list of dicts to headers + rows
                headers = list(data[0].keys()) if data else []
                rows = [[row.get(h) for h in headers] for row in data]
                sheets = [{"name": "Sheet1", "headers": headers, "rows": rows}]
            elif sheets is None:
                return ToolResult(success=False, error="Provide either `sheets` or `data`.")

            wb = Workbook()
            wb.remove(wb.active)  # remove default empty sheet

            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
            header_align = Alignment(horizontal="center", vertical="center")

            total_rows_written = 0

            for sheet_def in sheets:
                sheet_name = sheet_def.get("name", "Sheet")
                headers_row = sheet_def.get("headers", [])
                data_rows = sheet_def.get("rows", [])

                ws = wb.create_sheet(title=sheet_name[:31])  # Excel limit: 31 chars

                # Write headers
                for col_idx, header in enumerate(headers_row, start=1):
                    cell = ws.cell(row=1, column=col_idx, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_align

                # Write data rows
                for row_idx, row in enumerate(data_rows, start=2):
                    for col_idx, value in enumerate(row, start=1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
                    total_rows_written += 1

                # Auto-column width
                for col_idx, header in enumerate(headers_row, start=1):
                    col_letter = get_column_letter(col_idx)
                    max_width = len(str(header))
                    for row in data_rows:
                        val = row[col_idx - 1] if col_idx - 1 < len(row) else ""
                        max_width = max(max_width, len(str(val)) if val is not None else 0)
                    ws.column_dimensions[col_letter].width = min(max_width + 4, 60)

                # Freeze top row
                ws.freeze_panes = "A2"

                # Auto-filter
                if headers_row:
                    last_col = get_column_letter(len(headers_row))
                    ws.auto_filter.ref = f"A1:{last_col}1"

            wb.save(str(filepath))

            if self._state_manager is not None:
                await self._state_manager.register_file(filename, str(filepath))

            logger.info("xlsx_written", filepath=str(filepath), sheets=len(sheets))
            return ToolResult(
                success=True,
                data={"filepath": str(filepath), "sheets": len(sheets), "rows_written": total_rows_written},
            )

        except Exception as exc:
            logger.error("write_xlsx_failed", error=str(exc))
            return ToolResult(success=False, error=f"Failed to write XLSX: {exc}")


# ── WritePptxTool ─────────────────────────────────────────────────────────────

class WritePptxTool(BaseTool):
    """Create a PowerPoint presentation (.pptx)."""

    name = "write_pptx"
    description = (
        "Create a PowerPoint presentation (.pptx). "
        "Each slide can have a title, content (bullet points as list of strings or newline-separated text), "
        "layout ('title', 'content', or 'blank'), and optional speaker_notes. "
        "Returns the filepath."
    )
    parameters = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "Output filename (e.g. 'deck.pptx').",
            },
            "slides": {
                "type": "array",
                "description": "List of slide definitions.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {
                            "description": "Bullet points (list of strings or single string with newlines).",
                        },
                        "layout": {
                            "type": "string",
                            "enum": ["title", "content", "blank"],
                            "description": "Slide layout. Default: 'content' or 'title' for first slide.",
                        },
                        "speaker_notes": {
                            "type": "string",
                            "description": "Optional speaker notes for the slide.",
                        },
                    },
                    "required": ["title"],
                },
            },
            "output_dir": {
                "type": "string",
                "description": f"Directory to write the file. Defaults to '{_DEFAULT_OUTPUT_DIR}'.",
            },
        },
        "required": ["filename", "slides"],
    }

    def __init__(
        self,
        output_dir: str = _DEFAULT_OUTPUT_DIR,
        state_manager: "StateManager | None" = None,
    ) -> None:
        self._output_dir = output_dir
        self._state_manager = state_manager

    async def execute(
        self,
        filename: str,
        slides: list[dict],
        output_dir: str | None = None,
    ) -> ToolResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return ToolResult(success=False, error="python-pptx is not installed. Run: pip install python-pptx")

        try:
            effective_dir = output_dir if output_dir is not None else self._output_dir
            output_path = _ensure_output_dir(effective_dir)
            if not filename.endswith(".pptx"):
                filename += ".pptx"
            filepath = _safe_filepath(output_path, filename)
            if filepath is None:
                return ToolResult(success=False, error="Path traversal blocked")

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Layout indices in the default template:
            # 0 = Title Slide, 1 = Title and Content, 6 = Blank
            layout_map = {
                "title": prs.slide_layouts[0],
                "content": prs.slide_layouts[1],
                "blank": prs.slide_layouts[6],
            }

            for i, slide_def in enumerate(slides):
                raw_layout = slide_def.get("layout")
                if raw_layout is None:
                    raw_layout = "title" if i == 0 else "content"
                layout = layout_map.get(raw_layout, prs.slide_layouts[1])

                slide = prs.slides.add_slide(layout)

                # Set title
                slide_title = slide_def.get("title", "")
                if slide.shapes.title is not None:
                    slide.shapes.title.text = slide_title

                # Set content
                content = slide_def.get("content")
                if content is not None:
                    # Find the body placeholder (index 1 or any content placeholder)
                    body_ph = None
                    for ph in slide.placeholders:
                        if ph.placeholder_format.idx != 0:  # not the title
                            body_ph = ph
                            break

                    if body_ph is not None:
                        tf = body_ph.text_frame
                        tf.clear()
                        # Content can be a list of strings or a single string
                        if isinstance(content, list):
                            bullets = content
                        else:
                            bullets = str(content).splitlines()

                        for j, bullet in enumerate(bullets):
                            if j == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0

                # Speaker notes
                notes_text = slide_def.get("speaker_notes")
                if notes_text:
                    notes_slide = slide.notes_slide
                    tf = notes_slide.notes_text_frame
                    tf.text = notes_text

            prs.save(str(filepath))

            if self._state_manager is not None:
                await self._state_manager.register_file(filename, str(filepath))

            logger.info("pptx_written", filepath=str(filepath), slides=len(slides))
            return ToolResult(
                success=True,
                data={"filepath": str(filepath), "slide_count": len(slides)},
            )

        except Exception as exc:
            logger.error("write_pptx_failed", error=str(exc))
            return ToolResult(success=False, error=f"Failed to write PPTX: {exc}")


# ── PdfReaderTool ─────────────────────────────────────────────────────────────

class PdfReaderTool(BaseTool):
    """Extract text content from a PDF file."""

    name = "read_pdf"
    description = (
        "Extract text content from a PDF file. "
        "Returns the extracted text (up to max_pages pages). "
        "Handles encrypted and corrupt PDFs gracefully."
    )
    parameters = {
        "type": "object",
        "properties": {
            "filepath": {
                "type": "string",
                "description": "Path to the PDF file to read.",
            },
            "max_pages": {
                "type": "integer",
                "description": "Maximum number of pages to extract (default 20).",
                "default": 20,
            },
        },
        "required": ["filepath"],
    }

    def __init__(self, allowed_dirs: list[str] | None = None) -> None:
        self._allowed_dirs = allowed_dirs

    async def execute(self, filepath: str, max_pages: int = 20) -> ToolResult:
        path = Path(filepath).resolve()

        # Allowed-directory check (prevents arbitrary file read)
        if self._allowed_dirs is not None:
            if not any(path.is_relative_to(Path(d).resolve()) for d in self._allowed_dirs):
                return ToolResult(
                    success=False,
                    error=f"Access denied: {filepath!r} is not under an allowed directory",
                )

        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {filepath}")
        if not path.is_file():
            return ToolResult(success=False, error=f"Path is not a file: {filepath}")
        if path.suffix.lower() != ".pdf":
            return ToolResult(success=False, error=f"File does not appear to be a PDF: {filepath}")

        # Try pymupdf first, fall back to pypdf
        try:
            return await self._read_with_pymupdf(path, max_pages)
        except ImportError:
            pass

        try:
            return await self._read_with_pypdf(path, max_pages)
        except ImportError:
            return ToolResult(
                success=False,
                error="No PDF library available. Run: pip install pymupdf",
            )

    async def _read_with_pymupdf(self, path: Path, max_pages: int) -> ToolResult:
        import fitz  # pymupdf

        try:
            doc = fitz.open(str(path))
        except Exception as exc:
            return ToolResult(success=False, error=f"Failed to open PDF: {exc}")

        if doc.needs_pass:
            doc.close()
            return ToolResult(success=False, error="PDF is encrypted/password-protected")

        total_pages = doc.page_count
        pages_to_read = min(total_pages, max_pages)
        text_parts = []

        for page_num in range(pages_to_read):
            try:
                page = doc[page_num]
                text_parts.append(f"--- Page {page_num + 1} ---\n{page.get_text()}")
            except Exception as exc:
                text_parts.append(f"--- Page {page_num + 1} --- [Error: {exc}]")

        doc.close()
        full_text = "\n\n".join(text_parts)
        logger.info("pdf_read_pymupdf", filepath=str(path), pages=pages_to_read)
        return ToolResult(
            success=True,
            data={
                "text": full_text,
                "pages_read": pages_to_read,
                "total_pages": total_pages,
                "truncated": total_pages > max_pages,
            },
        )

    async def _read_with_pypdf(self, path: Path, max_pages: int) -> ToolResult:
        from pypdf import PdfReader

        try:
            reader = PdfReader(str(path))
        except Exception as exc:
            return ToolResult(success=False, error=f"Failed to open PDF: {exc}")

        if reader.is_encrypted:
            return ToolResult(success=False, error="PDF is encrypted/password-protected")

        total_pages = len(reader.pages)
        pages_to_read = min(total_pages, max_pages)
        text_parts = []

        for page_num in range(pages_to_read):
            try:
                text_parts.append(f"--- Page {page_num + 1} ---\n{reader.pages[page_num].extract_text()}")
            except Exception as exc:
                text_parts.append(f"--- Page {page_num + 1} --- [Error: {exc}]")

        full_text = "\n\n".join(text_parts)
        logger.info("pdf_read_pypdf", filepath=str(path), pages=pages_to_read)
        return ToolResult(
            success=True,
            data={
                "text": full_text,
                "pages_read": pages_to_read,
                "total_pages": total_pages,
                "truncated": total_pages > max_pages,
            },
        )
