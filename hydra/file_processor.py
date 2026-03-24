"""
FileProcessor — processes uploaded files, saves to output dir, extracts text.

Supports text extraction for common formats:
- Plain text / markup / config: .txt, .md, .csv, .json, .yaml, .yml, .xml, .html,
  .log, .ini, .toml, .cfg, .env
- Code: .py, .js, .ts, .java, .cpp, .c, .h, .go, .rs, .rb, .php, .swift, .kt
- Office: .pdf (pymupdf), .docx (python-docx), .xlsx (openpyxl), .pptx (python-pptx)

Unsupported formats (images, binaries, etc.) get extracted_text=None.
"""

from __future__ import annotations

import asyncio
import mimetypes
import shutil
import uuid
from pathlib import Path
from typing import TYPE_CHECKING

import structlog

from hydra.models import FileAttachment

if TYPE_CHECKING:
    pass

logger = structlog.get_logger(__name__)

_MAX_EXTRACTED_CHARS = 50_000
_TRUNCATION_MARKER = "\n[...truncated]"

# Extension → MIME type mapping (no magic bytes needed for v1)
_MIME_MAP: dict[str, str] = {
    # Text / markup
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    ".yaml": "application/yaml",
    ".yml": "application/yaml",
    ".xml": "application/xml",
    ".html": "text/html",
    ".htm": "text/html",
    # Config / log
    ".log": "text/plain",
    ".ini": "text/plain",
    ".toml": "application/toml",
    ".cfg": "text/plain",
    ".env": "text/plain",
    # Code
    ".py": "text/x-python",
    ".js": "application/javascript",
    ".ts": "application/typescript",
    ".java": "text/x-java",
    ".cpp": "text/x-c++",
    ".c": "text/x-c",
    ".h": "text/x-c",
    ".go": "text/x-go",
    ".rs": "text/x-rust",
    ".rb": "text/x-ruby",
    ".php": "application/x-php",
    ".swift": "text/x-swift",
    ".kt": "text/x-kotlin",
    # Office
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Extensions that support text extraction
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm",
    ".log", ".ini", ".toml", ".cfg", ".env",
    ".py", ".js", ".ts", ".java", ".cpp", ".c", ".h", ".go", ".rs", ".rb",
    ".php", ".swift", ".kt",
}
_PDF_EXTENSION = ".pdf"
_DOCX_EXTENSION = ".docx"
_XLSX_EXTENSION = ".xlsx"
_PPTX_EXTENSION = ".pptx"


def _detect_mime(filepath: Path) -> str | None:
    """Detect MIME type from file extension."""
    ext = filepath.suffix.lower()
    if ext in _MIME_MAP:
        return _MIME_MAP[ext]
    # Fall back to mimetypes stdlib
    mime, _ = mimetypes.guess_type(str(filepath))
    return mime


def _truncate(text: str) -> str:
    """Truncate text to MAX_EXTRACTED_CHARS, appending marker if truncated."""
    if len(text) <= _MAX_EXTRACTED_CHARS:
        return text
    return text[:_MAX_EXTRACTED_CHARS] + _TRUNCATION_MARKER


def _extract_text_sync(filepath: Path, mime_type: str | None) -> str | None:
    """
    Synchronous text extraction. Returns None for unsupported formats.
    Errors during extraction are caught and return None.
    """
    ext = filepath.suffix.lower()

    # Plain text / code / config files
    if ext in _TEXT_EXTENSIONS:
        try:
            text = filepath.read_text(encoding="utf-8", errors="replace")
            return _truncate(text)
        except Exception as e:
            logger.warning("text_extract_failed", filepath=str(filepath), error=str(e))
            return None

    # PDF via pymupdf
    if ext == _PDF_EXTENSION:
        try:
            import pymupdf  # type: ignore
            doc = pymupdf.open(str(filepath))
            parts: list[str] = []
            total_chars = 0
            for page in doc:
                page_text = page.get_text()
                parts.append(page_text)
                total_chars += len(page_text)
                if total_chars >= _MAX_EXTRACTED_CHARS:
                    break
            doc.close()
            return _truncate("".join(parts))
        except Exception as e:
            logger.warning("pdf_extract_failed", filepath=str(filepath), error=str(e))
            return None

    # DOCX via python-docx
    if ext == _DOCX_EXTENSION:
        try:
            from docx import Document  # type: ignore
            doc = Document(str(filepath))
            paragraphs = [p.text for p in doc.paragraphs]
            return _truncate("\n".join(paragraphs))
        except Exception as e:
            logger.warning("docx_extract_failed", filepath=str(filepath), error=str(e))
            return None

    # XLSX via openpyxl
    if ext == _XLSX_EXTENSION:
        try:
            import openpyxl  # type: ignore
            _XLSX_MAX_SHEETS = 20
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            all_sheets = wb.sheetnames
            sheets_to_process = all_sheets[:_XLSX_MAX_SHEETS]
            remaining_sheets = len(all_sheets) - len(sheets_to_process)
            parts: list[str] = []
            parts.append(f"Sheets: {', '.join(all_sheets)}")
            for sheet_name in sheets_to_process:
                ws = wb[sheet_name]
                parts.append(f"\n--- Sheet: {sheet_name} ---")
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if v is None else str(v) for v in row))
                    row_count += 1
                    if row_count >= 100:  # limit rows per sheet
                        parts.append("[... more rows truncated ...]")
                        break
            if remaining_sheets > 0:
                parts.append(f"[...{remaining_sheets} more sheets not shown]")
            wb.close()
            return _truncate("\n".join(parts))
        except Exception as e:
            logger.warning("xlsx_extract_failed", filepath=str(filepath), error=str(e))
            return None

    # PPTX via python-pptx
    if ext == _PPTX_EXTENSION:
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(filepath))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            return _truncate("\n\n".join(parts))
        except Exception as e:
            logger.warning("pptx_extract_failed", filepath=str(filepath), error=str(e))
            return None

    # Unsupported format
    return None


class FileProcessor:
    """
    Processes uploaded files: saves to output dir, extracts text where possible.

    Key methods:
    - process(files): takes a list of file paths (CLI/programmatic use)
    - process_upload(filename, content): takes raw bytes (API/frontend use)
    """

    def __init__(self, output_dir: str) -> None:
        self.output_dir = Path(output_dir)
        self.upload_dir = self.output_dir / "uploads"
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        logger.debug("file_processor_init", upload_dir=str(self.upload_dir))

    async def process(self, files: list[str | Path]) -> list[FileAttachment]:
        """
        Process a list of file paths. Returns FileAttachment objects with extracted text.

        Files that don't exist are handled gracefully (FileAttachment with error hint).
        """
        results: list[FileAttachment] = []
        for f in files:
            attachment = await self._process_single_path(Path(f))
            results.append(attachment)
        return results

    async def process_upload(self, filename: str, content: bytes) -> FileAttachment:
        """
        Process a raw upload (from API). Saves bytes to upload_dir, extracts text.
        """
        # Validate filename for null bytes and control characters
        if '\x00' in filename or any(ord(c) < 32 for c in filename):
            logger.warning("upload_invalid_filename", filename=repr(filename))
            return FileAttachment(
                filepath="",
                original_name=filename,
                mime_type=None,
                extracted_text=None,
            )

        try:
            # Make a unique destination to avoid collisions
            unique_stem = f"{uuid.uuid4().hex[:8]}_{Path(filename).stem}"
            dest_name = unique_stem + Path(filename).suffix
            dest_path = self.upload_dir / dest_name

            # Path traversal guard — ensure destination stays within upload_dir
            if not dest_path.resolve().is_relative_to(self.upload_dir.resolve()):
                raise ValueError(f"Unsafe upload path: {dest_name}")

            try:
                dest_path.write_bytes(content)
            except Exception as e:
                logger.error("upload_write_failed", filename=filename, error=str(e))
                return FileAttachment(
                    filepath=str(dest_path),
                    original_name=filename,
                    mime_type=_detect_mime(Path(filename)),
                    size_bytes=len(content),
                    extracted_text=None,
                )

            attachment = await self._process_single_path(dest_path, original_name=filename)
            return attachment
        except ValueError:
            raise
        except Exception as e:
            logger.error("process_upload_failed", filename=filename, error=str(e))
            return FileAttachment(
                filepath="",
                original_name=filename,
                mime_type=_detect_mime(Path(filename)),
                extracted_text=None,
            )

    def _extract_text(self, filepath: Path, mime_type: str | None) -> str | None:
        """Extract text from a file based on its type. Returns None if unsupported."""
        return _extract_text_sync(filepath, mime_type)

    # ── Private ───────────────────────────────────────────────────────────────

    async def _process_single_path(
        self, filepath: Path, original_name: str | None = None
    ) -> FileAttachment:
        """Process a single file path into a FileAttachment."""
        original_name = original_name or filepath.name

        if not filepath.exists():
            logger.warning("file_not_found", filepath=str(filepath))
            return FileAttachment(
                filepath=str(filepath.resolve()),
                original_name=original_name,
                mime_type=_detect_mime(filepath),
                size_bytes=0,
                extracted_text=None,
            )

        filepath = filepath.resolve()
        mime_type = _detect_mime(filepath)
        size_bytes = filepath.stat().st_size

        # Run text extraction in a thread to avoid blocking the event loop
        try:
            extracted_text = await asyncio.to_thread(
                _extract_text_sync, filepath, mime_type
            )
        except Exception as e:
            logger.warning("extract_failed", filepath=str(filepath), error=str(e))
            extracted_text = None

        logger.info(
            "file_processed",
            name=original_name,
            size_bytes=size_bytes,
            mime_type=mime_type,
            has_text=extracted_text is not None,
        )

        return FileAttachment(
            filepath=str(filepath),
            original_name=original_name,
            mime_type=mime_type,
            size_bytes=size_bytes,
            extracted_text=extracted_text,
        )
