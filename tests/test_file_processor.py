"""
Tests for FileProcessor and FileAttachment.
"""

from __future__ import annotations

import asyncio
import os
import tempfile
from pathlib import Path

import pytest
import pytest_asyncio

from hydra.file_processor import FileProcessor, _MAX_EXTRACTED_CHARS, _TRUNCATION_MARKER
from hydra.models import FileAttachment


# ── Fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture()
def tmp_dir(tmp_path):
    return tmp_path


@pytest.fixture()
def processor(tmp_dir):
    return FileProcessor(output_dir=str(tmp_dir / "output"))


# ── Helper ────────────────────────────────────────────────────────────────────

def make_file(directory: Path, name: str, content: str | bytes) -> Path:
    p = directory / name
    if isinstance(content, bytes):
        p.write_bytes(content)
    else:
        p.write_text(content, encoding="utf-8")
    return p


# ── MIME type detection ───────────────────────────────────────────────────────

def test_mime_type_txt(tmp_dir, processor):
    f = make_file(tmp_dir, "hello.txt", "hello world")
    asyncio.run(_check_mime(processor, f, "text/plain"))


def test_mime_type_py(tmp_dir, processor):
    f = make_file(tmp_dir, "script.py", "print('hi')")
    asyncio.run(_check_mime(processor, f, "text/x-python"))


def test_mime_type_pdf(tmp_dir, processor):
    # Just check MIME detection — actual PDF requires valid bytes
    from hydra.file_processor import _detect_mime
    assert _detect_mime(Path("report.pdf")) == "application/pdf"


def test_mime_type_unknown(tmp_dir, processor):
    from hydra.file_processor import _detect_mime
    # Unknown extension may return None or a guessed type — just ensure it doesn't crash
    result = _detect_mime(Path("archive.xyz"))
    # None or a string — both acceptable
    assert result is None or isinstance(result, str)


async def _check_mime(processor, filepath, expected_mime):
    att = await processor._process_single_path(filepath)
    assert att.mime_type == expected_mime


# ── Text extraction: plain text ───────────────────────────────────────────────

@pytest.mark.asyncio
async def test_txt_extraction(tmp_dir, processor):
    content = "Hello, this is a plain text file.\nSecond line."
    f = make_file(tmp_dir, "notes.txt", content)
    attachments = await processor.process([str(f)])
    assert len(attachments) == 1
    att = attachments[0]
    assert att.extracted_text == content
    assert att.original_name == "notes.txt"
    assert att.size_bytes > 0


@pytest.mark.asyncio
async def test_md_extraction(tmp_dir, processor):
    content = "# Title\n\nSome markdown content."
    f = make_file(tmp_dir, "readme.md", content)
    attachments = await processor.process([str(f)])
    assert attachments[0].extracted_text == content


@pytest.mark.asyncio
async def test_json_extraction(tmp_dir, processor):
    content = '{"key": "value", "number": 42}'
    f = make_file(tmp_dir, "data.json", content)
    attachments = await processor.process([str(f)])
    assert attachments[0].extracted_text == content


@pytest.mark.asyncio
async def test_py_extraction(tmp_dir, processor):
    content = "def hello():\n    return 'world'\n"
    f = make_file(tmp_dir, "code.py", content)
    attachments = await processor.process([str(f)])
    assert attachments[0].extracted_text == content


# ── Text extraction: docx ─────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_docx_extraction(tmp_dir, processor):
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not available")

    doc = Document()
    doc.add_paragraph("First paragraph.")
    doc.add_paragraph("Second paragraph.")
    docx_path = tmp_dir / "test.docx"
    doc.save(str(docx_path))

    attachments = await processor.process([str(docx_path)])
    assert attachments[0].extracted_text is not None
    assert "First paragraph." in attachments[0].extracted_text
    assert "Second paragraph." in attachments[0].extracted_text


# ── Text extraction: xlsx ─────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_xlsx_extraction(tmp_dir, processor):
    try:
        import openpyxl
    except ImportError:
        pytest.skip("openpyxl not available")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Month", "Amount"])
    ws.append(["Jan", 1000])
    ws.append(["Feb", 2000])
    xlsx_path = tmp_dir / "data.xlsx"
    wb.save(str(xlsx_path))

    attachments = await processor.process([str(xlsx_path)])
    text = attachments[0].extracted_text
    assert text is not None
    assert "Revenue" in text
    assert "Month" in text or "Jan" in text


# ── Text extraction: pptx ─────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_pptx_extraction(tmp_dir, processor):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not available")

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    slide.placeholders[1].text = "Some slide content."
    pptx_path = tmp_dir / "presentation.pptx"
    prs.save(str(pptx_path))

    attachments = await processor.process([str(pptx_path)])
    text = attachments[0].extracted_text
    assert text is not None
    assert "Test Slide" in text or "Some slide content." in text


# ── Text extraction: pdf ──────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_pdf_extraction(tmp_dir, processor):
    try:
        import pymupdf
    except ImportError:
        pytest.skip("pymupdf not available")

    # Create a minimal PDF with text
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF world!")
    pdf_path = tmp_dir / "test.pdf"
    doc.save(str(pdf_path))
    doc.close()

    attachments = await processor.process([str(pdf_path)])
    text = attachments[0].extracted_text
    assert text is not None
    assert "Hello PDF world!" in text


# ── Unsupported formats ───────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_binary_image_no_text(tmp_dir, processor):
    """PNG and other binary formats should return extracted_text=None."""
    # Write a fake PNG (not actually valid, but extension-based detection)
    f = make_file(tmp_dir, "photo.png", b"\x89PNG\r\n\x1a\n")
    attachments = await processor.process([str(f)])
    assert attachments[0].extracted_text is None
    assert attachments[0].size_bytes > 0


@pytest.mark.asyncio
async def test_mp3_no_text(tmp_dir, processor):
    f = make_file(tmp_dir, "audio.mp3", b"ID3\x03\x00\x00\x00\x00\x00\x00")
    attachments = await processor.process([str(f)])
    assert attachments[0].extracted_text is None


# ── Truncation ────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_truncation_at_50k(tmp_dir, processor):
    """Files with more than 50K chars should be truncated."""
    big_content = "A" * (_MAX_EXTRACTED_CHARS + 1000)
    f = make_file(tmp_dir, "big.txt", big_content)
    attachments = await processor.process([str(f)])
    text = attachments[0].extracted_text
    assert text is not None
    assert len(text) <= _MAX_EXTRACTED_CHARS + len(_TRUNCATION_MARKER)
    assert _TRUNCATION_MARKER in text


@pytest.mark.asyncio
async def test_no_truncation_under_50k(tmp_dir, processor):
    """Files at or below 50K chars should not be truncated."""
    content = "B" * (_MAX_EXTRACTED_CHARS - 100)
    f = make_file(tmp_dir, "medium.txt", content)
    attachments = await processor.process([str(f)])
    text = attachments[0].extracted_text
    assert text == content
    assert _TRUNCATION_MARKER not in text


# ── process_upload ────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_process_upload_saves_file(tmp_dir, processor):
    """process_upload should save the file and return a FileAttachment."""
    content = b"uploaded content here"
    att = await processor.process_upload("my_file.txt", content)

    assert att.original_name == "my_file.txt"
    assert att.size_bytes == len(content)
    # File should exist on disk
    assert Path(att.filepath).exists()
    assert Path(att.filepath).read_bytes() == content
    # Text should be extracted
    assert att.extracted_text == content.decode("utf-8")


@pytest.mark.asyncio
async def test_process_upload_binary(tmp_dir, processor):
    """process_upload for a binary file should save correctly with extracted_text=None."""
    content = b"\x00\x01\x02\x03\x04"
    att = await processor.process_upload("data.bin", content)
    assert att.original_name == "data.bin"
    assert Path(att.filepath).exists()
    assert att.extracted_text is None


@pytest.mark.asyncio
async def test_process_upload_unique_names(tmp_dir, processor):
    """Uploading two files with the same name should produce two distinct files."""
    att1 = await processor.process_upload("doc.txt", b"first")
    att2 = await processor.process_upload("doc.txt", b"second")
    assert att1.filepath != att2.filepath
    assert Path(att1.filepath).exists()
    assert Path(att2.filepath).exists()


# ── Non-existent file ─────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_nonexistent_file_graceful(tmp_dir, processor):
    """Non-existent files should not raise; return FileAttachment with size_bytes=0."""
    fake_path = str(tmp_dir / "does_not_exist.txt")
    attachments = await processor.process([fake_path])
    assert len(attachments) == 1
    att = attachments[0]
    assert att.size_bytes == 0
    assert att.extracted_text is None
    assert att.original_name == "does_not_exist.txt"


# ── Multiple files ────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_process_multiple_files(tmp_dir, processor):
    """process() should handle multiple files."""
    f1 = make_file(tmp_dir, "a.txt", "content A")
    f2 = make_file(tmp_dir, "b.py", "# content B")
    f3 = make_file(tmp_dir, "c.png", b"\x89PNG")

    attachments = await processor.process([str(f1), str(f2), str(f3)])
    assert len(attachments) == 3
    assert attachments[0].extracted_text == "content A"
    assert attachments[1].extracted_text == "# content B"
    assert attachments[2].extracted_text is None
