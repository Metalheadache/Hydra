"""
Tests for document tools, HTTP request tool, and data transform tool.
"""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest


# ── WriteDocxTool ─────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_write_docx_creates_valid_file():
    from hydra_agents.tools.document_tools import WriteDocxTool
    from docx import Document

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteDocxTool(output_dir=tmpdir)
        content = (
            "# Introduction\n\n"
            "This is a **bold** statement.\n\n"
            "## Details\n\n"
            "- First item\n"
            "- Second item\n\n"
            "Plain paragraph here."
        )
        result = await tool.execute(filename="test.docx", content=content, title="Test Document")
        assert result.success, f"Expected success, got error: {result.error}"
        filepath = result.data["filepath"]
        assert Path(filepath).exists()

        # Validate it's a real docx
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        assert len(paragraphs) > 0, "Document should have content"


@pytest.mark.asyncio
async def test_write_docx_adds_extension():
    from hydra_agents.tools.document_tools import WriteDocxTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteDocxTool(output_dir=tmpdir)
        result = await tool.execute(filename="noext", content="Hello world")
        assert result.success
        assert result.data["filepath"].endswith(".docx")


@pytest.mark.asyncio
async def test_write_docx_path_traversal_blocked():
    from hydra_agents.tools.document_tools import WriteDocxTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteDocxTool(output_dir=tmpdir)
        result = await tool.execute(filename="../evil.docx", content="pwned")
        assert not result.success
        assert "traversal" in result.error.lower()


# ── WriteXlsxTool ─────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_write_xlsx_creates_valid_file():
    from hydra_agents.tools.document_tools import WriteXlsxTool
    from openpyxl import load_workbook

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteXlsxTool(output_dir=tmpdir)
        sheets = [
            {
                "name": "Sales",
                "headers": ["Month", "Revenue", "Units"],
                "rows": [
                    ["January", 10000, 150],
                    ["February", 12000, 180],
                    ["March", 9500, 140],
                ],
            }
        ]
        result = await tool.execute(filename="sales.xlsx", sheets=sheets)
        assert result.success, f"Expected success, got: {result.error}"
        assert result.data["rows_written"] == 3

        filepath = result.data["filepath"]
        wb = load_workbook(filepath)
        ws = wb["Sales"]
        assert ws.cell(1, 1).value == "Month"
        assert ws.cell(2, 1).value == "January"
        assert ws.cell(4, 2).value == 9500


@pytest.mark.asyncio
async def test_write_xlsx_from_data_dicts():
    from hydra_agents.tools.document_tools import WriteXlsxTool
    from openpyxl import load_workbook

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteXlsxTool(output_dir=tmpdir)
        data = [
            {"name": "Alice", "score": 95},
            {"name": "Bob", "score": 87},
        ]
        result = await tool.execute(filename="scores.xlsx", data=data)
        assert result.success
        wb = load_workbook(result.data["filepath"])
        ws = wb.active or wb.worksheets[0]
        assert ws.cell(1, 1).value == "name"
        assert ws.cell(2, 1).value == "Alice"


@pytest.mark.asyncio
async def test_write_xlsx_no_data_returns_error():
    from hydra_agents.tools.document_tools import WriteXlsxTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WriteXlsxTool(output_dir=tmpdir)
        result = await tool.execute(filename="empty.xlsx")
        assert not result.success


# ── WritePptxTool ─────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_write_pptx_creates_valid_file():
    from hydra_agents.tools.document_tools import WritePptxTool
    from pptx import Presentation

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WritePptxTool(output_dir=tmpdir)
        slides = [
            {"title": "Title Slide", "content": "Subtitle here", "layout": "title"},
            {"title": "Slide 2", "content": ["Point A", "Point B", "Point C"], "layout": "content"},
            {"title": "Slide 3", "content": "Summary", "speaker_notes": "Speak about X"},
        ]
        result = await tool.execute(filename="deck.pptx", slides=slides)
        assert result.success, f"Expected success, got: {result.error}"
        assert result.data["slide_count"] == 3

        prs = Presentation(result.data["filepath"])
        assert len(prs.slides) == 3


@pytest.mark.asyncio
async def test_write_pptx_correct_slide_count():
    from hydra_agents.tools.document_tools import WritePptxTool
    from pptx import Presentation

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = WritePptxTool(output_dir=tmpdir)
        slides = [{"title": f"Slide {i}"} for i in range(5)]
        result = await tool.execute(filename="five.pptx", slides=slides)
        assert result.success
        prs = Presentation(result.data["filepath"])
        assert len(prs.slides) == 5


# ── PdfReaderTool ─────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_pdf_reader_handles_missing_file():
    from hydra_agents.tools.document_tools import PdfReaderTool

    tool = PdfReaderTool()
    result = await tool.execute(filepath="/nonexistent/path/file.pdf")
    assert not result.success
    assert "not found" in result.error.lower()


@pytest.mark.asyncio
async def test_pdf_reader_rejects_non_pdf():
    from hydra_agents.tools.document_tools import PdfReaderTool

    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        f.write(b"just a text file")
        tmppath = f.name

    tool = PdfReaderTool()
    result = await tool.execute(filepath=tmppath)
    assert not result.success
    Path(tmppath).unlink(missing_ok=True)


@pytest.mark.asyncio
async def test_pdf_reader_reads_real_pdf():
    """Create a real PDF with pymupdf and read it back."""
    try:
        import fitz
    except ImportError:
        pytest.skip("pymupdf not available")

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = Path(tmpdir) / "sample.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello from Hydra PDF test!")
        doc.save(str(pdf_path))
        doc.close()

        from hydra_agents.tools.document_tools import PdfReaderTool
        tool = PdfReaderTool()
        result = await tool.execute(filepath=str(pdf_path))
        assert result.success, f"Expected success, got: {result.error}"
        assert "Hello from Hydra" in result.data["text"]
        assert result.data["pages_read"] == 1


# ── HttpRequestTool ───────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_http_request_blocks_localhost():
    from hydra_agents.tools.research_tools import HttpRequestTool

    tool = HttpRequestTool()
    for url in [
        "http://localhost/",
        "http://127.0.0.1/",
        "http://127.0.0.1:8080/api",
    ]:
        result = await tool.execute(method="GET", url=url)
        assert not result.success, f"Should have blocked SSRF for {url}"
        assert "ssrf" in result.error.lower() or "blocked" in result.error.lower()


@pytest.mark.asyncio
async def test_http_request_blocks_private_networks():
    from hydra_agents.tools.research_tools import HttpRequestTool

    tool = HttpRequestTool()
    for url in [
        "http://10.0.0.1/",
        "http://192.168.1.1/",
        "http://172.16.0.1/",
        "http://172.31.255.255/",
    ]:
        result = await tool.execute(method="GET", url=url)
        assert not result.success, f"Should have blocked SSRF for {url}"


@pytest.mark.asyncio
async def test_http_request_rejects_non_http_scheme():
    from hydra_agents.tools.research_tools import HttpRequestTool

    tool = HttpRequestTool()
    result = await tool.execute(method="GET", url="ftp://example.com/file")
    assert not result.success
    assert "http" in result.error.lower()


@pytest.mark.asyncio
async def test_http_request_rejects_invalid_method():
    from hydra_agents.tools.research_tools import HttpRequestTool

    tool = HttpRequestTool()
    result = await tool.execute(method="INVALID", url="https://example.com")
    assert not result.success


# ── DataTransformTool ─────────────────────────────────────────────────────────

SAMPLE_DATA = [
    {"name": "Alice", "dept": "Eng", "salary": 90000, "age": 30},
    {"name": "Bob", "dept": "Eng", "salary": 80000, "age": 25},
    {"name": "Carol", "dept": "HR", "salary": 70000, "age": 35},
    {"name": "Dave", "dept": "HR", "salary": 65000, "age": 28},
    {"name": "Eve", "dept": "Eng", "salary": 95000, "age": 32},
]


@pytest.mark.asyncio
async def test_data_transform_filter():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "filter", "params": {"field": "dept", "operator": "==", "value": "Eng"}}],
    )
    assert result.success
    rows = result.data["result"]
    assert len(rows) == 3
    assert all(r["dept"] == "Eng" for r in rows)


@pytest.mark.asyncio
async def test_data_transform_filter_numeric():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "filter", "params": {"field": "salary", "operator": ">=", "value": 80000}}],
    )
    assert result.success
    rows = result.data["result"]
    assert len(rows) == 3
    assert all(r["salary"] >= 80000 for r in rows)


@pytest.mark.asyncio
async def test_data_transform_sort_asc():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "sort", "params": {"field": "salary", "order": "asc"}}],
    )
    assert result.success
    rows = result.data["result"]
    salaries = [r["salary"] for r in rows]
    assert salaries == sorted(salaries)


@pytest.mark.asyncio
async def test_data_transform_sort_desc():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "sort", "params": {"field": "salary", "order": "desc"}}],
    )
    assert result.success
    rows = result.data["result"]
    salaries = [r["salary"] for r in rows]
    assert salaries == sorted(salaries, reverse=True)


@pytest.mark.asyncio
async def test_data_transform_group_by_count():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "group_by", "params": {"field": "dept", "agg_func": "count"}}],
    )
    assert result.success
    rows = result.data["result"]
    counts = {r["dept"]: r["count"] for r in rows}
    assert counts["Eng"] == 3
    assert counts["HR"] == 2


@pytest.mark.asyncio
async def test_data_transform_group_by_sum():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[
            {
                "type": "group_by",
                "params": {"field": "dept", "agg_field": "salary", "agg_func": "sum"},
            }
        ],
    )
    assert result.success
    rows = result.data["result"]
    sums = {r["dept"]: r["sum_salary"] for r in rows}
    assert sums["Eng"] == 90000 + 80000 + 95000
    assert sums["HR"] == 70000 + 65000


@pytest.mark.asyncio
async def test_data_transform_select():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "select", "params": {"fields": ["name", "salary"]}}],
    )
    assert result.success
    rows = result.data["result"]
    assert all(set(r.keys()) == {"name", "salary"} for r in rows)


@pytest.mark.asyncio
async def test_data_transform_limit():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "limit", "params": {"count": 2}}],
    )
    assert result.success
    assert len(result.data["result"]) == 2


@pytest.mark.asyncio
async def test_data_transform_pipeline():
    """Filter → sort → select → limit chain."""
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[
            {"type": "filter", "params": {"field": "dept", "operator": "==", "value": "Eng"}},
            {"type": "sort", "params": {"field": "salary", "order": "desc"}},
            {"type": "select", "params": {"fields": ["name", "salary"]}},
            {"type": "limit", "params": {"count": 2}},
        ],
    )
    assert result.success
    rows = result.data["result"]
    assert len(rows) == 2
    assert rows[0]["salary"] > rows[1]["salary"]
    assert "dept" not in rows[0]


@pytest.mark.asyncio
async def test_data_transform_unknown_op_returns_error():
    from hydra_agents.tools.data_tools import DataTransformTool

    tool = DataTransformTool()
    result = await tool.execute(
        data=SAMPLE_DATA,
        operations=[{"type": "explode", "params": {}}],
    )
    assert not result.success


# ── ChartGeneratorTool ────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_chart_generator_bar():
    from hydra_agents.tools.data_tools import ChartGeneratorTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = ChartGeneratorTool(output_dir=tmpdir)
        result = await tool.execute(
            chart_type="bar",
            data={"labels": ["A", "B", "C"], "values": [10, 20, 15]},
            title="Bar Chart",
            filename="bar.png",
        )
        assert result.success, f"Expected success, got: {result.error}"
        assert Path(result.data["filepath"]).exists()


@pytest.mark.asyncio
async def test_chart_generator_pie():
    from hydra_agents.tools.data_tools import ChartGeneratorTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = ChartGeneratorTool(output_dir=tmpdir)
        result = await tool.execute(
            chart_type="pie",
            data={"labels": ["X", "Y"], "values": [60, 40]},
            title="Pie Chart",
            filename="pie.png",
        )
        assert result.success
        assert Path(result.data["filepath"]).exists()


@pytest.mark.asyncio
async def test_chart_generator_invalid_type():
    from hydra_agents.tools.data_tools import ChartGeneratorTool

    with tempfile.TemporaryDirectory() as tmpdir:
        tool = ChartGeneratorTool(output_dir=tmpdir)
        result = await tool.execute(
            chart_type="heatmap",
            data={"labels": ["A"], "values": [1]},
            title="Bad",
            filename="bad.png",
        )
        assert not result.success


# ── ToolRegistry integration ──────────────────────────────────────────────────

def test_registry_registers_new_tools():
    from hydra_agents.tool_registry import ToolRegistry

    registry = ToolRegistry()
    registry.register_defaults()

    expected_tools = [
        "write_docx",
        "write_xlsx",
        "write_pptx",
        "read_pdf",
        "generate_chart",
        "data_transform",
        "http_request",
        "translate_text",
        "summarize_text",
    ]
    for name in expected_tools:
        assert name in registry, f"Tool '{name}' not found in registry"


def test_registry_total_tool_count():
    from hydra_agents.tool_registry import ToolRegistry

    registry = ToolRegistry()
    registry.register_defaults()
    # Should have at least 22 tools (13 original + 9 new)
    assert len(registry) >= 22, f"Expected ≥22 tools, got {len(registry)}"
